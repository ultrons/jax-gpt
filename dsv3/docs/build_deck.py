"""
Build DSv3 MoE Performance Analysis deck.
Run:  uv run --with python-pptx python3 build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x0F, 0x0F)   # near-black
ACCENT    = RGBColor(0x00, 0xC8, 0xFF)   # bright cyan
ACCENT2   = RGBColor(0xFF, 0x6B, 0x35)   # orange
GREEN     = RGBColor(0x39, 0xD3, 0x53)   # green
YELLOW    = RGBColor(0xFF, 0xD6, 0x00)   # yellow
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x88, 0x88, 0x88)
DARKGREY  = RGBColor(0x22, 0x22, 0x22)
CODEFG    = RGBColor(0xD4, 0xD4, 0xD4)
CODEBG    = RGBColor(0x1E, 0x1E, 0x1E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl

def box(sl, x, y, w, h, text, size=18, bold=False, color=WHITE,
        bg=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    return txb

def hline(sl, y, x0=0.3, x1=13.0, color=ACCENT, width=Pt(1)):
    from pptx.util import Inches
    line = sl.shapes.add_connector(1, Inches(x0), Inches(y), Inches(x1), Inches(y))
    line.line.color.rgb = color
    line.line.width = width

def rect(sl, x, y, w, h, fill=DARKGREY, line_color=None, line_w=Pt(1)):
    from pptx.util import Inches
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def code_box(sl, x, y, w, h, code, size=13):
    """Monospaced dark code block."""
    rect(sl, x, y, w, h, fill=CODEBG)
    box(sl, x+0.12, y+0.08, w-0.24, h-0.16, code,
        size=size, color=CODEFG, bg=None)

def title_line(sl, text, sub=None):
    box(sl, 0.4, 0.2, 12.5, 0.7, text, size=32, bold=True, color=WHITE)
    hline(sl, 1.0)
    if sub:
        box(sl, 0.4, 1.05, 12.5, 0.45, sub, size=16, color=GREY, italic=True)

def bullet(sl, x, y, items, size=17, color=WHITE, indent=0.25, spacing=0.38):
    for i, item in enumerate(items):
        box(sl, x, y + i*spacing, 11.5, 0.4,
            ("  " * indent) + "• " + item, size=size, color=color)

def highlight_box(sl, x, y, w, h, label, value, unit="",
                  label_color=GREY, val_color=ACCENT):
    rect(sl, x, y, w, h, fill=DARKGREY, line_color=ACCENT, line_w=Pt(1.5))
    box(sl, x+0.1, y+0.08, w-0.2, 0.3, label, size=13, color=label_color,
        align=PP_ALIGN.CENTER)
    box(sl, x+0.1, y+0.35, w-0.2, 0.45, value, size=26, bold=True,
        color=val_color, align=PP_ALIGN.CENTER)
    if unit:
        box(sl, x+0.1, y+0.78, w-0.2, 0.25, unit, size=12, color=GREY,
            align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 1.5, 1.8, 10.5, 1.2,
    "DSv3 MoE Training Performance",
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(sl, 1.5, 3.1, 10.5, 0.6,
    "TPU v7x  ·  4×8×8  ·  512 devices",
    size=22, color=ACCENT, align=PP_ALIGN.CENTER)
box(sl, 1.5, 3.8, 10.5, 0.5,
    "BS=1024  ·  EP=8  ·  FSDP=64",
    size=18, color=GREY, align=PP_ALIGN.CENTER)
hline(sl, 5.0, 2.0, 11.3, color=ACCENT)
box(sl, 1.5, 5.15, 10.5, 0.5,
    "Roofline analysis  ·  ICI bottleneck breakdown  ·  Engineering roadmap",
    size=15, color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 2 — v7x hardware
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "v7x Hardware: What We're Working With")

# hardware spec boxes
specs = [
    ("MXU\n(bf16)", "2,307", "TFLOP/s per chip"),
    ("HBM\nBW", "7,373", "GB/s per chip"),
    ("ICI\n(cross-chip)", "~90", "GB/s per link"),
    ("ICI\n(same chip)", "600", "GB/s intra-chip"),
    ("HBM\ncapacity", "95", "GB per chip"),
]
for i, (lbl, val, unit) in enumerate(specs):
    highlight_box(sl, 0.35 + i*2.55, 1.5, 2.3, 1.15, lbl, val, unit)

box(sl, 0.4, 2.85, 12.5, 0.35,
    "2 cores (JAX devices) per chip  →  per-device: 1,153 TFLOP/s · 3,686 GB/s HBM · 32 GB VMEM",
    size=14, color=GREY)

hline(sl, 3.35, color=DARKGREY)

box(sl, 0.4, 3.45, 12.5, 0.4,
    "Roofline: when is an op compute-bound vs memory-bound?", size=18, bold=True, color=ACCENT)

code = (
    "# Operational intensity (OI) = FLOPs / bytes_accessed\n"
    "# Ridge point = peak_compute / peak_memory_bw\n"
    "\n"
    "peak_flops_bf16 = 1153e12   # FLOP/s per device\n"
    "peak_hbm_bw    = 3686e9    # bytes/s per device\n"
    "\n"
    "ridge_point = peak_flops_bf16 / peak_hbm_bw  # OI threshold\n"
    "# → 1153e12 / 3686e9 = 313 FLOP/byte\n"
    "\n"
    "# If OI > 313: compute-bound  (MXU is the bottleneck)\n"
    "# If OI < 313: memory-bound   (HBM bandwidth is the bottleneck)"
)
code_box(sl, 0.4, 3.9, 8.0, 3.1, code, size=13)

box(sl, 8.7, 3.9, 4.3, 0.4, "Ridge point", size=15, bold=True, color=ACCENT)
box(sl, 8.7, 4.3, 4.3, 0.4, "313 FLOP/byte", size=28, bold=True, color=YELLOW,
    align=PP_ALIGN.LEFT)
box(sl, 8.7, 4.8, 4.3, 1.6,
    "An op needs to do 313\nFLOPs per byte read/written\nto saturate the MXU.\n\nBelow that — you're just\nwaiting on HBM.",
    size=15, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 3 — Expert FFN OI
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "Expert FFN: Is It Compute-Bound?",
           sub="Operational intensity of the core matmuls — before any communication")

code = (
    "# Per-expert forward pass (one of 32 local experts)\n"
    "# tpe = 1,536 tokens  D = 7,168  F = 2,048\n"
    "\n"
    "h_gate = tok_e @ W_gate.T   # (tpe, D) × (D, F)  →  (tpe, F)\n"
    "h_up   = tok_e @ W_up.T    # (tpe, D) × (D, F)  →  (tpe, F)\n"
    "act    = silu(h_gate) * h_up\n"
    "out_e  = act @ W_out.T     # (tpe, F) × (F, D)  →  (tpe, D)\n"
    "\n"
    "# FLOPs (one expert, forward):\n"
    "flops_w1 = 2 * tpe * D * F   # = 2 × 1536 × 7168 × 2048 = 45 GFLOP\n"
    "flops_w2 = 2 * tpe * F * D   # = 2 × 1536 × 2048 × 7168 = 45 GFLOP  (same)\n"
    "# ×2 because gate+up share the same D→F projection\n"
    "total_fwd = 2 * (flops_w1 + flops_w2) = 181 GFLOP  per expert\n"
    "\n"
    "# Bytes read (weights, bf16, once per expert):\n"
    "bytes_W_gate = D * F * 2 = 29.4 MB\n"
    "bytes_W_up   = D * F * 2 = 29.4 MB\n"
    "bytes_W_out  = F * D * 2 = 29.4 MB\n"
    "total_bytes  = 88 MB  (plus 88 MB tokens + activations)"
)
code_box(sl, 0.4, 1.25, 7.8, 5.5, code, size=12)

# OI callout
rect(sl, 8.5, 1.25, 4.4, 2.3, fill=DARKGREY, line_color=ACCENT, line_w=Pt(2))
box(sl, 8.6, 1.3, 4.2, 0.4, "OI — expert FFN", size=14, bold=True, color=ACCENT)
box(sl, 8.6, 1.7, 4.2, 0.5,
    "181 GFLOP / 176 MB\n= 1,028 FLOP/byte",
    size=18, bold=True, color=GREEN)
box(sl, 8.6, 2.3, 4.2, 0.6,
    "Ridge point: 313 FLOP/byte\n→ 3× above ridge",
    size=14, color=WHITE)
box(sl, 8.6, 2.85, 4.2, 0.5,
    "✓ Expert FFN is compute-bound\n   (MXU is the bottleneck for FFN)",
    size=13, color=GREEN)

rect(sl, 8.5, 3.75, 4.4, 3.0, fill=DARKGREY, line_color=ACCENT2, line_w=Pt(2))
box(sl, 8.6, 3.8, 4.2, 0.4, "But: 32 experts × 15 ms = ?", size=14, bold=True, color=ACCENT2)
box(sl, 8.6, 4.25, 4.2, 0.5,
    "Compute per device\nper MoE layer:",
    size=14, color=WHITE)
box(sl, 8.6, 4.75, 4.2, 0.5,
    "17,391 GFLOP (4× fwd)\n@ 1,153 TFLOP/s",
    size=15, bold=True, color=YELLOW)
box(sl, 8.6, 5.3, 4.2, 0.4,
    "= 15 ms per MoE layer",
    size=18, bold=True, color=GREEN)
box(sl, 8.6, 5.75, 4.2, 0.4,
    "Problem: ICI takes 28 ms.\nCompute is not the limit.",
    size=13, color=ACCENT2)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 4 — Token flow / config
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "Token Flow: From GBS to tpe",
           sub="How 4M training tokens get distributed to 32 local experts per device")

code = (
    "# Global config\n"
    "GBS     = 1024           # sequences in this batch\n"
    "seqlen  = 4096           # tokens per sequence\n"
    "n_total_tokens = GBS * seqlen   # 4,194,304\n"
    "\n"
    "# Device count\n"
    "EP, FSDP = 8, 64\n"
    "n_devices = EP * FSDP   # 512  (256 chips × 2 cores)\n"
    "\n"
    "# Per-device token shard  (act_spec = P(('ep','fsdp'), None))\n"
    "T_local = n_total_tokens // n_devices   # 8,192 tokens\n"
    "\n"
    "# Inside shard_map: EP all_gather reconstructs FSDP-local view\n"
    "T_fsdp  = T_local * EP                  # 65,536 tokens\n"
    "# All 8 EP ranks within a FSDP group see the same 65,536 tokens\n"
    "\n"
    "# Routing  (DSv3: 256 routed experts, top-6)\n"
    "K, E = 6, 256\n"
    "E_local = E // EP                        # 32 experts per EP rank\n"
    "\n"
    "token_expert_pairs = T_fsdp * K         # 393,216\n"
    "pairs_for_this_rank = token_expert_pairs // EP   # 49,152\n"
    "tpe = pairs_for_this_rank // E_local    # 1,536  (avg)\n"
    "\n"
    "# + 2 shared experts (replicated, no A2A needed)"
)
code_box(sl, 0.4, 1.2, 7.6, 5.9, code, size=13)

# flow diagram (text-art boxes)
box(sl, 8.35, 1.3, 4.5, 0.35, "Token pipeline", size=14, bold=True, color=ACCENT)

steps = [
    ("4,194,304 tokens", "global batch", WHITE, DARKGREY),
    ("8,192 tokens", "T_local / device", WHITE, DARKGREY),
    ("65,536 tokens", "T_fsdp after EP gather", YELLOW, DARKGREY),
    ("49,152 pairs", "for this rank's 32 experts", GREEN, DARKGREY),
    ("1,536 avg", "tokens per expert (tpe)", ACCENT, DARKGREY),
]
for i, (val, lbl, vc, bc) in enumerate(steps):
    y = 1.75 + i * 1.05
    rect(sl, 8.35, y, 4.5, 0.85, fill=bc, line_color=ACCENT, line_w=Pt(1))
    box(sl, 8.45, y+0.04, 4.3, 0.38, val, size=18, bold=True, color=vc)
    box(sl, 8.45, y+0.43, 4.3, 0.3, lbl, size=12, color=GREY)
    if i < len(steps)-1:
        box(sl, 10.3, y+0.88, 0.5, 0.17, "↓", size=14, color=GREY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 5 — EP topology + ICI breakdown
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "EP Topology: Not All ICI Is Equal",
           sub="EP=8 spans 4 chips × 2 cores. Same-chip core pairs communicate at 600 GB/s.")

# chip diagram
box(sl, 0.4, 1.2, 12.5, 0.38,
    "Physical layout: EP=8 on 4 chips, 2 cores each",
    size=15, bold=True, color=ACCENT)

chip_colors = [RGBColor(0x1A,0x3A,0x5C), RGBColor(0x1A,0x3A,0x5C),
               RGBColor(0x1A,0x3A,0x5C), RGBColor(0x1A,0x3A,0x5C)]
for c in range(4):
    cx = 0.5 + c * 3.1
    rect(sl, cx, 1.7, 2.8, 1.5, fill=chip_colors[c],
         line_color=ACCENT, line_w=Pt(2))
    box(sl, cx+0.05, 1.72, 2.7, 0.32, f"Chip {c}", size=13, bold=True, color=ACCENT)
    # core 0
    rect(sl, cx+0.1, 2.1, 1.2, 0.85, fill=RGBColor(0x0A,0x2A,0x4A),
         line_color=GREEN, line_w=Pt(1))
    box(sl, cx+0.12, 2.15, 1.15, 0.3, f"EP rank {c*2}", size=11, bold=True, color=GREEN)
    box(sl, cx+0.12, 2.45, 1.15, 0.3, "core 0", size=10, color=GREY)
    # core 1
    rect(sl, cx+1.5, 2.1, 1.2, 0.85, fill=RGBColor(0x0A,0x2A,0x4A),
         line_color=GREEN, line_w=Pt(1))
    box(sl, cx+1.52, 2.15, 1.15, 0.3, f"EP rank {c*2+1}", size=11, bold=True, color=GREEN)
    box(sl, cx+1.52, 2.45, 1.15, 0.3, "core 1", size=10, color=GREY)

# intra-chip arrow
for c in range(4):
    cx = 0.5 + c * 3.1
    box(sl, cx+1.3, 2.38, 0.25, 0.3, "⟺", size=16, color=YELLOW, align=PP_ALIGN.CENTER)
box(sl, 4.5, 3.3, 4.3, 0.3, "600 GB/s (intra-chip)", size=13, color=YELLOW,
    align=PP_ALIGN.CENTER)

# inter-chip arrows
for c in range(3):
    cx = 0.5 + c * 3.1 + 2.8
    box(sl, cx+0.05, 2.35, 0.25, 0.3, "→", size=16, color=ACCENT2, align=PP_ALIGN.CENTER)
box(sl, 0.4, 3.55, 12.5, 0.3, "~90 GB/s per ICI link (cross-chip)", size=13, color=ACCENT2,
    align=PP_ALIGN.CENTER)

hline(sl, 4.0, color=DARKGREY)

# traffic table
box(sl, 0.4, 4.1, 12.5, 0.38,
    "A2A traffic per device per MoE layer  (forward token send, K=6, T_fsdp=65,536)",
    size=14, bold=True, color=WHITE)

rows = [
    ("Destination",           "Pairs",    "Volume",    "Bandwidth",   "Time",    WHITE,   True),
    ("Local EP rank (own)",   "49,152",   "0 bytes",   "—",           "0 ms",    GREEN,   False),
    ("Co-chip EP rank (×1)",  "49,152",   "703 MB",    "600 GB/s",    "~1 ms",   YELLOW,  False),
    ("Cross-chip EP (×6)",    "49,152 ea","703 MB ea", "~90 GB/s",    "~8 ms",   ACCENT2, False),
    ("Total ICI (slow path)", "295,000",  "4.2 GB",    "150 GB/s agg","~28 ms",  ACCENT,  False),
]
col_x = [0.4, 3.1, 5.3, 7.3, 9.7, 11.3]
col_w = [2.6, 2.1, 2.0, 2.3, 1.5, 1.5]
for r, (c1,c2,c3,c4,c5,clr,hdr) in enumerate(rows):
    y = 4.55 + r * 0.46
    if hdr:
        rect(sl, 0.35, y-0.04, 12.6, 0.42, fill=RGBColor(0x1A,0x1A,0x2E), line_color=None)
    for ci, (txt, cx, cw) in enumerate(zip([c1,c2,c3,c4,c5], col_x, col_w)):
        box(sl, cx, y, cw, 0.38, txt, size=13,
            color=(GREY if hdr else clr) if ci > 0 else (WHITE if hdr else clr),
            bold=hdr)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 6 — Why ICI dominates
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "The Bottleneck: ICI, Not MXU",
           sub="Even at peak MXU efficiency, A2A ICI is 2× slower than expert FFN compute")

code = (
    "# Per MoE layer, per device\n"
    "\n"
    "# --- Compute ---\n"
    "flops_per_layer = 32 * 1536 * (2*7168*4096 + 2*2048*7168) * 4  # fwd+bwd+recompute\n"
    "#                = 17,391 GFLOP\n"
    "compute_time_ms = 17_391 / 1_153_000  * 1000   # @ 1153 TFLOP/s\n"
    "#               = 15 ms\n"
    "\n"
    "# --- ICI  (cross-chip only; local + co-chip are free) ---\n"
    "tokens_for_remote_experts = T_fsdp * K * 6/8   # 6 cross-chip EP ranks\n"
    "#                         = 65536 * 6 * 0.75 = 295,000 token slots\n"
    "bytes_one_direction = 295_000 * 7168 * 2        # BF16\n"
    "#                   = 4.2 GB\n"
    "ici_time_ms = 4.2e9 / 150e9 * 1000              # 150 GB/s aggregate\n"
    "#           = 28 ms\n"
    "\n"
    "# 4 transfers on critical path (fwd send, fwd recv, bwd send, bwd recv)\n"
    "# With perfect layer-to-layer pipelining → effective ~56 ms / layer\n"
    "# Without pipelining                     → 4 × 28 = 112 ms / layer\n"
    "\n"
    "# ICI / Compute ratio  =  28 ms / 15 ms  =  1.87×\n"
    "# Even MXU at 100% can't hide the ICI cost without pipelining."
)
code_box(sl, 0.4, 1.2, 7.8, 5.5, code, size=12)

# bar chart (text-art)
box(sl, 8.5, 1.2, 4.5, 0.4, "Per-layer time budget", size=14, bold=True, color=ACCENT)
bars = [
    ("Compute (MXU)",    15,  GREEN,   "15 ms"),
    ("ICI one direction",28,  ACCENT2, "28 ms"),
    ("ICI 4× (no pipe)", 112, RGBColor(0xCC,0x22,0x22), "112 ms"),
    ("ICI (pipelined)",  56,  YELLOW,  "56 ms"),
]
max_ms = 120
bar_w = 3.5
for i, (lbl, ms, clr, txt) in enumerate(bars):
    y = 1.8 + i * 1.1
    w = bar_w * ms / max_ms
    box(sl, 8.5, y, 4.5, 0.3, lbl, size=12, color=WHITE)
    rect(sl, 8.5, y+0.32, w, 0.45, fill=clr)
    box(sl, 8.5+w+0.05, y+0.32, 1.2, 0.45, txt, size=13, bold=True, color=clr)

hline(sl, 6.45, color=DARKGREY)
box(sl, 0.4, 6.5, 12.5, 0.4,
    "Key insight: SC offload frees MXU from scatter/gather, but MXU is not the bottleneck — ICI is. "
    "The real gain is from A2A DMA pipelining across layers.",
    size=13, color=GREY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 7 — A2A pipelining
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "A2A Pipelining: The Critical Optimisation",
           sub="Overlap ICI of layer N with expert FFN compute of layer N−1")

code = (
    "# Naive (sequential):  every layer stalls on ICI\n"
    "# ─────────────────────────────────────────────────\n"
    "for layer in moe_layers:\n"
    "    tokens = a2a_send(sorted_tokens)   # 28 ms — STALL\n"
    "    out    = expert_ffn(tokens)        # 15 ms\n"
    "    result = a2a_recv(out)             # 28 ms — STALL\n"
    "# per layer: 28+15+28 = 71 ms  →  58 layers = 4.1 s\n"
    "\n"
    "# Pipelined (DMA runs concurrently with compute):\n"
    "# ─────────────────────────────────────────────────\n"
    "issue_a2a(layer[0].tokens)             # start layer 0 ICI\n"
    "for i, layer in enumerate(moe_layers):\n"
    "    wait_a2a(layer[i].tokens)          # wait for this layer's tokens\n"
    "    if i+1 < len(moe_layers):\n"
    "        issue_a2a(layer[i+1].tokens)   # prefetch next layer's tokens\n"
    "    out = expert_ffn(layer[i].tokens)  # compute overlaps with ^ ICI\n"
    "    issue_a2a(layer[i].out)            # send results back\n"
    "    wait_a2a(layer[i-1].out)           # collect previous results\n"
    "# per layer: max(28, 15) = 28 ms  →  58 layers = 1.6 s\n"
    "# (bwd further overlaps: effective ~56 ms fwd+bwd / layer)"
)
code_box(sl, 0.4, 1.25, 8.0, 5.5, code, size=12)

# pipeline diagram
box(sl, 8.5, 1.25, 4.5, 0.4, "ICI / Compute timeline", size=14, bold=True, color=ACCENT)

tscale = 0.025   # inches per ms
y0 = 1.8
layers_shown = 3
for li in range(layers_shown):
    lbl = f"Layer {li}"
    row_y = y0 + li * 1.55

    box(sl, 8.5, row_y, 0.8, 0.3, lbl, size=11, color=GREY)

    # ICI send bar (28 ms)
    x_ici_send = 8.5 + li * 0.7
    rect(sl, x_ici_send, row_y+0.35, 28*tscale, 0.38, fill=ACCENT2,
         line_color=None)
    box(sl, x_ici_send, row_y+0.35, 28*tscale, 0.38,
        "ICI↑", size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # Compute bar (15 ms, overlapping with ICI)
    x_comp = x_ici_send + 28*tscale - 15*tscale*0.5
    rect(sl, x_comp, row_y+0.78, 15*tscale, 0.38, fill=GREEN, line_color=None)
    box(sl, x_comp, row_y+0.78, 15*tscale, 0.38,
        "FFN", size=9, color=BG, align=PP_ALIGN.CENTER, bold=True)

    # ICI recv bar
    x_ici_recv = x_comp + 15*tscale
    rect(sl, x_ici_recv, row_y+1.2, 28*tscale, 0.38, fill=ACCENT,
         line_color=None)
    box(sl, x_ici_recv, row_y+1.2, 28*tscale, 0.38,
        "ICI↓", size=9, color=BG, align=PP_ALIGN.CENTER, bold=True)

box(sl, 8.5, y0+layers_shown*1.55+0.1, 4.5, 0.5,
    "ICI↑ and ICI↓ of adjacent layers\noverlap → 28 ms critical path",
    size=12, color=YELLOW)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 8 — Step time estimate
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "Step Time Estimate: Full Optimised Stack",
           sub="BS=1024 · EP=8 · FSDP=64 · 58 MoE + 61 MLA + 3 dense · fwd+bwd+recompute")

# component budget
box(sl, 0.4, 1.15, 12.5, 0.38,
    "Per-step time budget  (with perfect A2A pipeline + FSDP overlap)", size=15, bold=True, color=WHITE)

components = [
    ("58 MoE layers  (ICI-bound, ~56 ms/layer pipelined)",  "3.2 s", ACCENT2, 3.2),
    ("61 MLA attention layers  (compute-bound)",             "0.13 s", GREEN,  0.13),
    ("3 dense MLP layers",                                   "0.02 s", GREEN,  0.02),
    ("FSDP weight all-gather  (overlapped w/ compute)",      "~0",     GREY,   0.0),
    ("TOTAL",                                                "~3.4 s", YELLOW, 3.35),
]
max_bar = 3.5
bar_area = 5.0
for i, (lbl, val, clr, secs) in enumerate(components):
    y = 1.65 + i * 0.8
    bw = bar_area * secs / max_bar if secs > 0 else 0.05
    rect(sl, 0.4, y+0.08, bw, 0.45, fill=clr if clr != GREY else DARKGREY)
    box(sl, 0.5, y+0.1, bw+0.1, 0.38, "", size=1)   # spacer
    box(sl, 0.4+bw+0.12, y+0.08, 6.5-bw, 0.45, lbl, size=14, color=WHITE)
    box(sl, 9.9, y+0.08, 2.8, 0.45, val, size=16, bold=True, color=clr,
        align=PP_ALIGN.RIGHT)
    if i == len(components)-2:
        hline(sl, y+0.65, x0=0.3, x1=13.0, color=GREY, width=Pt(0.75))

hline(sl, 5.4, color=DARKGREY)

# scenario table
box(sl, 0.4, 5.5, 12.5, 0.38,
    "Scenarios", size=15, bold=True, color=WHITE)

rows = [
    ("Scenario",                        "Step time", "TPS/chip", "vs v10",  WHITE,  True),
    ("Optimistic  (perfect pipeline)",  "3.4 s",     "~4,800",   "+10.3×", GREEN,  False),
    ("Realistic   (60% efficiency)",    "5.5 s",     "~2,960",   "+6.4×",  YELLOW, False),
    ("Conservative (40% efficiency)",   "8 s",       "~2,040",   "+4.4×",  WHITE,  False),
    ("v10 ragged_dot  (baseline)",      "~35 s",     "465",      "1.0×",   GREY,   False),
]
col_x2 = [0.4, 4.2, 6.8, 9.0, 11.2]
col_w2 = [3.7, 2.5, 2.1, 2.1, 1.8]
for r, row in enumerate(rows):
    *cells, clr, hdr = row
    y = 5.95 + r * 0.44
    if hdr:
        rect(sl, 0.35, y-0.04, 12.6, 0.42, fill=RGBColor(0x1A,0x1A,0x2E))
    for ci, (txt, cx, cw) in enumerate(zip(cells, col_x2, col_w2)):
        box(sl, cx, y, cw, 0.38, txt, size=13,
            color=(GREY if hdr else clr),
            bold=(hdr or (ci == 3 and not hdr)))

# ═══════════════════════════════════════════════════════════════════════════
# Slide 9 — What eliminates the EP all-reduce
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "Why A2A Eliminates the EP All-Reduce",
           sub="Current streaming_bwd uses psum(d_tok, 'ep') per layer — the most expensive collective")

code_left = (
    "# Current (streaming_bwd v1): EP all-reduce\n"
    "# ─────────────────────────────────────────\n"
    "def _streaming_bwd_fn(g_l, fx_l, fw_l, fi_l,\n"
    "                      w0_l, w1_l, wout_l):\n"
    "    # EP all_gather to see all T/FSDP tokens\n"
    "    fx_full = lax.all_gather(fx_l, 'ep')  # ICI\n"
    "\n"
    "    # ... expert FFN backward ...\n"
    "    d_tok_partial = per_expert_backward()\n"
    "\n"
    "    # ALL-REDUCE: every EP rank needs to sum\n"
    "    # contributions from ALL other EP ranks\n"
    "    d_tok_full = lax.psum(d_tok_partial, 'ep')\n"
    "    #  ^^ 65536 × 7168 × 4B × 2 (all-reduce)\n"
    "    #   = 3.76 GB per layer per device\n"
    "    #   = 3.76 / 150 GB/s = 25 ms per layer\n"
    "    #   = 25 ms × 58 layers = 1.45 s extra\n"
    "    # Plus: global sync point — no overlap"
)
code_box(sl, 0.4, 1.2, 6.2, 4.8, code_left, size=11.5)

code_right = (
    "# Pallas A2A: point-to-point DMA, no psum\n"
    "# ─────────────────────────────────────────\n"
    "# Forward:\n"
    "#   sort tokens → DMA sorted_tokens to\n"
    "#   the 7 EP peers that own their experts\n"
    "#   → each expert rank receives exactly\n"
    "#     the tokens it needs. No broadcast.\n"
    "\n"
    "# Backward:\n"
    "#   each expert rank computes d_tok_e\n"
    "#   DMA d_tok back to token-owner ranks\n"
    "#   token owners scatter-add contributions\n"
    "#   → no all-reduce needed at all\n"
    "\n"
    "# ICI traffic comparison per layer:\n"
    "#   all-reduce:  3.76 GB  (blocking)\n"
    "#   A2A send:    4.2 GB   (async DMA)\n"
    "#   A2A recv:    4.2 GB   (async DMA)\n"
    "# Total A2A = more bytes but ASYNC →\n"
    "# overlaps with compute of next layer."
)
code_box(sl, 6.8, 1.2, 6.1, 4.8, code_right, size=11.5)

# cost comparison
rect(sl, 0.4, 6.15, 5.8, 0.95, fill=RGBColor(0x33,0x11,0x11), line_color=RGBColor(0xCC,0x22,0x22))
box(sl, 0.5, 6.2, 5.6, 0.4, "Current EP psum: 58 × 25 ms = 1.45 s  (blocking)", size=13, color=ACCENT2, bold=True)
box(sl, 0.5, 6.6, 5.6, 0.4, "Global sync every layer — cannot pipeline", size=12, color=GREY)

rect(sl, 6.8, 6.15, 6.1, 0.95, fill=RGBColor(0x11,0x33,0x11), line_color=GREEN)
box(sl, 6.9, 6.2, 5.9, 0.4, "A2A DMA: async, pipelines with FFN compute", size=13, color=GREEN, bold=True)
box(sl, 6.9, 6.6, 5.9, 0.4, "No blocking collective — ICI hidden in compute shadow", size=12, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 10 — Engineering roadmap
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "Engineering Roadmap",
           sub="Four work items to go from current streaming_bwd to full Pallas A2A kernel")

items = [
    {
        "n": "1",
        "title": "Pallas Forward Kernel  (Expert FFN A2A)",
        "status": "NOT STARTED",
        "status_color": ACCENT2,
        "detail": [
            "Replace shard_map+psum EP forward with Pallas A2A",
            "sort tokens → DMA send → per-expert FFN → DMA recv → scatter-add",
            "SC-offloaded indexed gather for token packing (sc_gather_rows)",
            "Validate: loss matches jax baseline at full scale",
        ],
        "impact": "Eliminates EP all_gather in fwd (65536×7168 per layer)",
    },
    {
        "n": "2",
        "title": "Pallas Backward Kernel  (Expert FFN Backward + A2A)",
        "status": "IN PROGRESS",
        "status_color": YELLOW,
        "detail": [
            "streaming_bwd v1/v2: JAX backward running, SC Get fix in v73",
            "Replace EP psum(d_tok) with reverse A2A DMA (no all-reduce)",
            "FSDP weight all-gather double-buffered across E_local experts",
            "Validate: loss curve 86.703→86.561 matching v60 baseline",
        ],
        "impact": "Eliminates EP psum (1.45 s/step) — biggest single win",
    },
    {
        "n": "3",
        "title": "Layer-to-Layer A2A Pipeline",
        "status": "NOT STARTED",
        "status_color": ACCENT2,
        "detail": [
            "Issue DMA for layer n+1 tokens while computing layer n FFN",
            "Double-buffer: two A2A slots (inflight + being processed)",
            "Requires ICI DMA API that returns before transfer completes",
            "Target: ICI fully hidden behind compute (56 ms → 15 ms/layer)",
        ],
        "impact": "Closes compute/ICI gap: ~3.4 s vs ~5.5 s step time",
    },
    {
        "n": "4",
        "title": "FSDP All-Gather Overlap  (Attention + Dense)",
        "status": "PARTIAL",
        "status_color": YELLOW,
        "detail": [
            "Async collective flags already in LIBTPU_INIT_ARGS (--xla_enable_async_all_gather)",
            "Per-layer all_gather issued before layer compute, reduce-scatter after",
            "Attention: 4 weight projections × 205 MB per layer × 61 layers = 50 GB",
            "Currently: XLA handles via async collectives; Pallas gives explicit control",
        ],
        "impact": "Hides 50 GB attention + 328 GB MoE weight ICI in compute shadow",
    },
]

for i, item in enumerate(items):
    col = i % 2
    row = i // 2
    x = 0.35 + col * 6.55
    y = 1.2 + row * 2.95

    rect(sl, x, y, 6.3, 2.75, fill=DARKGREY,
         line_color=item["status_color"], line_w=Pt(2))

    # header row
    rect(sl, x, y, 6.3, 0.48, fill=RGBColor(0x1A,0x1A,0x2E))
    box(sl, x+0.1, y+0.04, 0.35, 0.38, item["n"], size=18, bold=True, color=item["status_color"])
    box(sl, x+0.5, y+0.04, 4.3, 0.38, item["title"], size=13, bold=True, color=WHITE)
    box(sl, x+4.85, y+0.06, 1.35, 0.32, item["status"], size=10, bold=True,
        color=item["status_color"], align=PP_ALIGN.RIGHT)

    for j, d in enumerate(item["detail"]):
        box(sl, x+0.15, y+0.55+j*0.46, 6.0, 0.42, "• " + d, size=11.5, color=WHITE)

    rect(sl, x, y+2.38, 6.3, 0.37, fill=RGBColor(0x0A,0x0A,0x1A))
    box(sl, x+0.12, y+2.42, 6.0, 0.3,
        "Impact: " + item["impact"], size=11, color=ACCENT, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 11 — Summary
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
title_line(sl, "Summary: Where Performance Comes From")

takeaways = [
    (ACCENT,  "v7x ridge point is 313 FLOP/byte",
               "Expert FFN OI ≈ 1,028 — solidly compute-bound in isolation."),
    (GREEN,   "Expert FFN is NOT the bottleneck",
               "At 1,536 tpe, each expert takes 15 ms compute vs 28 ms ICI. ICI wins."),
    (YELLOW,  "EP A2A topology matters",
               "1/8 traffic is local (free). 1/8 is 600 GB/s co-chip (near-free). 6/8 is cross-chip at ~90 GB/s."),
    (ACCENT2, "EP all-reduce is the single biggest waste",
               "Current psum: 1.45 s/step, blocking. A2A DMA: async, pipelines with compute."),
    (GREEN,   "A2A layer pipelining is the key lever",
               "Perfect pipelining → 3.4 s step. Realistic → 5.5 s. Either is 6–10× over v10 (465 TPS/chip)."),
    (WHITE,   "SC offload helps but is not critical",
               "Frees MXU from scatter/gather, but MXU is already not the bottleneck."),
]

for i, (clr, heading, detail) in enumerate(takeaways):
    y = 1.15 + i * 1.0
    rect(sl, 0.35, y, 0.06, 0.55, fill=clr)
    box(sl, 0.55, y, 12.0, 0.38, heading, size=16, bold=True, color=clr)
    box(sl, 0.55, y+0.38, 12.0, 0.42, detail, size=13, color=GREY)

hline(sl, 7.1, color=ACCENT)
box(sl, 0.4, 7.18, 12.5, 0.28,
    "Realistic target with full Pallas A2A + pipelining:  ~5.5 s/step  ·  ~2,960 TPS/chip  ·  6× over v10",
    size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────
out = "/home/sivaibhav_google_com/ml-experiments/dsv3/docs/moe_performance_analysis.pptx"
prs.save(out)
print(f"Saved: {out}")
