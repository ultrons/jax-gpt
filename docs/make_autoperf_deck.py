"""Build a brief PPTX deck on the autoperf project for a senior-audience overview."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Theme ────────────────────────────────────────────────────────────────────
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xC0, 0x39, 0x4A)
GRAY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xEA, 0xEC, 0xEE)
GREEN = RGBColor(0x2C, 0x7A, 0x4B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=None,
             align=None, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align if align else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=GRAY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        # support (text, indent) tuples or plain strings
        if isinstance(item, tuple):
            text, indent = item
        else:
            text, indent = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = indent
        bullet = "• " if indent == 0 else "‒ "
        r = p.add_run()
        r.text = bullet + text
        r.font.name = "Calibri"
        r.font.size = Pt(size - indent * 2)
        r.font.color.rgb = color
    return tb

def add_title(slide, text, *, color=BLUE, size=32):
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.8),
             text, size=size, bold=True, color=color)
    # underline rule
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(1.05), Inches(1.5), Emu(45000))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def add_footer(slide, text):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
             text, size=10, color=GRAY)

def add_box(slide, x, y, w, h, *, fill=LIGHT, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    return shp

# ── Slide 1: title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_box(s, Emu(0), Emu(0), SW, SH, fill=BLUE)
add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.2),
         "Autoperf",
         size=72, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.6),
         "Agent-driven performance optimization for large-scale LLM training",
         size=24, color=RGBColor(0xCF, 0xD8, 0xDC))
add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
         "DSv3 671B · TPU v7x · 16 iterations · 2 cluster-validated wins",
         size=18, color=RGBColor(0xCF, 0xD8, 0xDC))
add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
         "Internal walkthrough · 2026-05-11",
         size=14, color=RGBColor(0xB0, 0xBE, 0xC5))

# ── Slide 2: framing ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "What is autoperf trying to do?")
add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
         "Optimize one training workload at a time, autonomously, on a real cluster.",
         size=20, color=GRAY)
add_bullets(s, Inches(0.7), Inches(2.3), Inches(12.0), Inches(4.0), [
    ("Pick the slowest leaf in the training step (per-leaf headroom report)", 0),
    ("Apply ONE focused change (kernel swap, sharding tweak, calibration fix)", 0),
    ("Build, run on TPU v7x, capture profile, measure", 0),
    ("Compare to the prior iteration; either commit or revert", 0),
    ("Repeat until top-leaf headroom is below threshold OR a halt condition fires", 0),
])
# emphasis box
box = add_box(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.3),
              fill=RGBColor(0xFF, 0xF4, 0xE5), border=ACCENT)
add_text(s, Inches(0.95), Inches(5.75), Inches(11.5), Inches(1.0), [
    "Why agents and not a Python orchestrator?",
    "We tried that path early — over-engineered, brittle. The prompt-as-spec, shell-as-tool",
    "shape (Karpathy-style autoresearch) gives us a system that can reason about cluster",
    "failures, write GitHub issues, and adjust strategy across iterations.",
], size=14, color=GRAY)
add_footer(s, "autoperf · framing")

# ── Slide 3: components — jax-gpt ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Component 1 — jax-gpt")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "The training repo. Where models live, where the iteration loop runs.",
         size=18, color=GRAY)
add_bullets(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.0), [
    ("Multi-model JAX/TPU repo (DSv3 671B, Qwen3.5 397B, GPT-2)", 0),
    ("Owns the full training pipeline: model code, sharding, optimizer, data, profile capture", 0),
    ("autoperf/ subdirectory hosts the agent harness:", 0),
    ("AGENT.md = the agent's role spec (~350 lines, system-prompt style)", 1),
    ("workloads/*.yaml = per-workload spec (model, hardware, sharding, overrides)", 1),
    ("v7x_KNOWLEDGE.md = anti-hallucination ledger (known bugs, trusted leaves)", 1),
    ("BLOCKED.md = ledger of open tool issues + autoperf-side tasks", 1),
    ("Production baseline: 1882 TPS/chip @ 30.5% MFU on v7x 4×8×8 (iter-2b confirmed)", 0),
    ("IMPROVED candidate: 1916 TPS/chip @ 31.1% MFU (iter-16, pending ratchet)", 0),
])
add_text(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.4),
         "github.com/ultrons/jax-gpt   ·   branch: autoperf/dsv3_train_full",
         size=12, color=BLUE)
add_footer(s, "component intro · 1 of 4")

# ── Slide 4: components — perfsim ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Component 2 — perfsim")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "The cost model. Roofline-with-microbench-calibration: tells you the floor.",
         size=18, color=GRAY)
add_bullets(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.0), [
    ("Predicts per-leaf time on a workload: \"Expert_gmm should take 13.7M μs/step\"", 0),
    ("Headroom = measured (cluster) − predicted (perfsim). That's the gap to close.", 0),
    ("Calibrated against real microbench data (GEMM efficiency curves per dtype/shape)", 0),
    ("Outputs a structured JSON headroom report with top-3 leaves ranked", 0),
    ("Boundary discipline: does NOT model kernel choices (gmm_v2 vs ragged_dot)", 0),
    ("That's the engineer's job; perfsim's job is just to say where the gap is", 1),
    ("Has experimental graph-based bandwidth estimator (topology module)", 0),
    ("for evaluating alternative network topologies (boardfly, twisted torus)", 1),
])
add_text(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.4),
         "github.com/ultrons/perfsim   ·   ~/perfsim → ~/ml-experiments-perfsim",
         size=12, color=BLUE)
add_footer(s, "component intro · 2 of 4")

# ── Slide 5: components — cde ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Component 3 — cde")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "The job manager. Build → run → profile → status, all behind one CLI.",
         size=18, color=GRAY)
add_bullets(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.0), [
    ("Hides the kubectl + JobSet + Docker complexity from agents and humans alike", 0),
    ("cde build = build trainer image, push to GCR (auto-tagged from build context hash)", 0),
    ("cde run --tag <id> --context <ctx> --profile = submit JobSet, capture xplane profile", 0),
    ("cde status / logs / profile / history = poll, stream, retrieve, audit", 0),
    ("Tracks every run in ~/.cde/history.sqlite (per-user run history)", 0),
    ("Agents call it as a tool; the agent never touches kubectl directly", 0),
    ("CI gate: pytest + mypy on every push and PR", 0),
])
add_text(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.4),
         "github.com/ultrons/cde   ·   src-layout, pip install -e .",
         size=12, color=BLUE)
add_footer(s, "component intro · 3 of 4")

# ── Slide 6: components — xla-shell ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Component 4 — xla-shell")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "The profile inspector. Surfaces raw data from XLA xplane.pb files.",
         size=18, color=GRAY)
add_bullets(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.0), [
    ("xprof-backed CLI for inspecting XLA profiles after a run", 0),
    ("Lists HLO ops, fusion records, hop times, multi-chip collective traces, roofline", 0),
    ("Surfaces RAW data — does NOT bucket fusion-records into perfsim's leaf names", 0),
    ("Bucketing is perfsim's domain (LEAF_PATTERNS in inference/scripts/headroom_report.py)", 1),
    ("This boundary keeps xla-shell stable as a parsing tool", 1),
    ("Provides get_op_shape API: given an op name + field, return tensor shape from HLO", 0),
    ("Enables defense-in-depth checks (compare perfsim's assumed dims vs xplane's actual)", 1),
    ("Real .xplane.pb fixtures in tests/ — synthetic tests miss xprof's binary surprises", 0),
])
add_text(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.4),
         "github.com/ultrons/xla-shell   ·   xprof pin >=2.21.6,<2.22",
         size=12, color=BLUE)
add_footer(s, "component intro · 4 of 4")

# ── Slide 7: architecture v1 — 4-agent ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Architecture v1 — 4 cooperating agents")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "One doer per repo, GitHub issues as the message bus.",
         size=18, color=GRAY)

# Diagram: autoperf in middle, 3 maintainer agents around it
def draw_box(s, x, y, w, h, label, sub, fill=LIGHT, txt=BLUE):
    box = add_box(s, x, y, w, h, fill=fill, border=BLUE)
    add_text(s, x, y + Inches(0.25), w, Inches(0.4), label,
             size=18, bold=True, color=txt, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.75), w, Inches(0.4), sub,
             size=12, color=GRAY, align=PP_ALIGN.CENTER)

# autoperf center
draw_box(s, Inches(5.0), Inches(2.4), Inches(3.3), Inches(1.4),
         "autoperf agent", "in jax-gpt", fill=RGBColor(0xFF, 0xE0, 0xB2))

# 3 maintainers
draw_box(s, Inches(0.7), Inches(4.8), Inches(3.0), Inches(1.3),
         "perfsim agent", "fixes calibration / model bugs")
draw_box(s, Inches(5.15), Inches(4.8), Inches(3.0), Inches(1.3),
         "cde agent", "fixes job-manager bugs")
draw_box(s, Inches(9.6), Inches(4.8), Inches(3.0), Inches(1.3),
         "xla-shell agent", "fixes parsing / schema bugs")

# arrows from autoperf to maintainers
def line(s, x1, y1, x2, y2):
    line = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = GRAY
    line.line.width = Pt(1.5)

line(s, 6.0, 3.8, 2.2, 4.8)
line(s, 6.65, 3.8, 6.65, 4.8)
line(s, 7.3, 3.8, 11.1, 4.8)

add_text(s, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.4),
         "Communication: GitHub issues. Each maintainer polls for autoperf-blocking; fixes; comments closed.",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
         "Audit trail comes free. Async. LLM-agnostic.",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, "architecture v1")

# ── Slide 8: small win — iter-2 ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "First win — iter-2: gmm_v2 enable")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "Single-line change. Real measured gain. Confirmed via headroom report.",
         size=18, color=GRAY)
# left: results table
from pptx.util import Inches as I
tbl = s.shapes.add_table(rows=5, cols=3,
                         left=Inches(0.7), top=Inches(2.3),
                         width=Inches(6.0), height=Inches(2.5)).table
hdr = ["Metric", "v304 baseline", "iter-2"]
rows = [
    ["step time", "37.0 s", "34.65 s (−2.35 s)"],
    ["TPS/chip", "1,770", "1,882 (+6.6%)"],
    ["MFU", "28.6%", "30.5% (+1.9 pp)"],
    ["step-1 loss", "415.491", "415.46–415.47 ✓"],
]
for j, txt in enumerate(hdr):
    cell = tbl.cell(0, j)
    cell.text = txt
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
for i, row in enumerate(rows, 1):
    for j, txt in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = txt
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)

# right: narrative
add_bullets(s, Inches(7.0), Inches(2.3), Inches(5.8), Inches(4.0), [
    ("Lever picked from heuristic table:", 0),
    ("\"verify gmm_ag kernel registration over default ragged-dot\"", 1),
    ("Targets the Expert_gmm leaf (top headroom)", 1),
    ("Single workload-yaml change (one boolean flag)", 0),
    ("Plus one inline import-bug fix found mid-iteration:", 0),
    ("model.py:1793: kernels.gmm_v2_train (relative import)", 1),
    ("This was the trigger to rethink the architecture", 1),
])

add_footer(s, "first win · 2026-05-07 · branch autoperf/dsv3_train_full")

# ── Slide 9: architecture v2 — pivot ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Architecture v2 — 1 doer, 3 reviewers")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "Why we pivoted: cross-repo handoffs added latency for fixes the doer could land in 30 sec.",
         size=18, color=GRAY)

# Diagram: autoperf with worktrees
draw_box(s, Inches(5.0), Inches(2.2), Inches(3.3), Inches(1.4),
         "autoperf agent", "fix-inline authority", fill=RGBColor(0xFF, 0xE0, 0xB2))

# Worktrees
add_text(s, Inches(0.7), Inches(4.0), Inches(12.0), Inches(0.4),
         "Per-repo worktrees on `autoperf-loop` branch (PYTHONPATH-isolated):",
         size=14, color=GRAY)
draw_box(s, Inches(0.7), Inches(4.5), Inches(3.0), Inches(0.9),
         "perfsim worktree", "~/autoperf/repos/perfsim/")
draw_box(s, Inches(5.15), Inches(4.5), Inches(3.0), Inches(0.9),
         "cde worktree", "~/autoperf/repos/cde/")
draw_box(s, Inches(9.6), Inches(4.5), Inches(3.0), Inches(0.9),
         "xla-shell worktree", "~/autoperf/repos/xla-shell/")

# Reviewers
add_text(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(0.4),
         "Hourly PR review on autoperf-loop branches (humans gate merges to main):",
         size=14, color=GRAY)
draw_box(s, Inches(0.7), Inches(6.05), Inches(3.0), Inches(0.85),
         "perfsim reviewer", "PR review only", fill=RGBColor(0xE3, 0xF2, 0xFD))
draw_box(s, Inches(5.15), Inches(6.05), Inches(3.0), Inches(0.85),
         "cde reviewer", "PR review only", fill=RGBColor(0xE3, 0xF2, 0xFD))
draw_box(s, Inches(9.6), Inches(6.05), Inches(3.0), Inches(0.85),
         "xla-shell reviewer", "PR review only", fill=RGBColor(0xE3, 0xF2, 0xFD))

# arrows
line(s, 6.0, 3.6, 2.2, 4.5)
line(s, 6.65, 3.6, 6.65, 4.5)
line(s, 7.3, 3.6, 11.1, 4.5)

add_footer(s, "architecture v2 · pivoted same-day after iter-2 inline-fix")

# ── Slide 10: validation — perfsim deep review ───────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Validation — perfsim deep review")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "13 issues filed at 04:32 UTC → 12 resolved by 06:02 UTC, same day.",
         size=18, color=GRAY)

add_bullets(s, Inches(0.7), Inches(2.3), Inches(6.0), Inches(4.0), [
    ("Triggered by autoperf's iter-2 caveat (bucketer staleness post-gmm_v2)", 0),
    ("Deep review surfaced 12 gaps in perfsim:", 0),
    ("ADR-001 was half-applied (training path on legacy code)", 1),
    ("Per-leaf confidence not propagated", 1),
    ("BF16 had no efficiency curve (scalar fallback)", 1),
    ("Search not yet wired to CLI", 1),
    ("Maintainer (reviewer) shipped 6 fixes in PR #20", 0),
    ("then 4 more in PR #22 (per-op port + thin swap)", 0),
    ("Net: 12 of 13 closed; only #10 (BF16 curve) remains, blocked on", 0),
    ("an autoperf-side microbench grid", 1),
])

# right: trust state restoration table
add_text(s, Inches(7.2), Inches(2.3), Inches(5.6), Inches(0.4),
         "Trust restoration on v304 leaves:",
         size=14, bold=True, color=BLUE)
tbl = s.shapes.add_table(rows=6, cols=3,
                         left=Inches(7.2), top=Inches(2.7),
                         width=Inches(5.6), height=Inches(2.6)).table
hdr = ["leaf", "pre", "post"]
rows = [
    ["Expert_gmm", "1.12", "1.18"],
    ["Attn_scores", "0.69", "0.82"],
    ["O_proj", "0.75", "0.92"],
    ["QKV_proj", "0.40", "0.48"],
    ["EP_AG_dispatch", "0.94", "0.99"],
]
for j, txt in enumerate(hdr):
    cell = tbl.cell(0, j)
    cell.text = txt
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(13)
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
for i, row in enumerate(rows, 1):
    for j, txt in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = txt
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12)

add_text(s, Inches(7.2), Inches(5.5), Inches(5.6), Inches(0.4),
         "Top-3 ranking shifted:",
         size=12, bold=True, color=BLUE)
add_text(s, Inches(7.2), Inches(5.85), Inches(5.6), Inches(0.4),
         "[FSDP_AG, Router, Norms]  →  [Expert_gmm, Norms, FSDP_AG]",
         size=12, color=GRAY)
add_text(s, Inches(7.2), Inches(6.25), Inches(5.6), Inches(0.4),
         "Expert_gmm now has positive headroom — what the iter loop actually wanted.",
         size=12, color=GREEN)

add_footer(s, "validation · maintainer agent's reviewer role landed in real time")

# ── Slide 11: iters 3-14 honest middle ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Iters 3–14 — the honest middle")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "4 days, 4 cluster shots, 0 perf wins on the iter-2b baseline. Worth doing anyway.",
         size=18, color=GRAY)

# Left column: cluster shots
add_text(s, Inches(0.5), Inches(2.1), Inches(6.2), Inches(0.4),
         "Cluster shots: all reverted",
         size=16, bold=True, color=BLUE)
add_bullets(s, Inches(0.5), Inches(2.55), Inches(6.2), Inches(3.0), [
    ("iter-5: tgmm tile_m=4096 — −1.4% TPS (memory-bound)", 0),
    ("iter-7: attn_proj_out OFFLOAD — NaN, filed jax-gpt#2", 0),
    ("iter-8: prevent_cse=True — NaN, filed jax-gpt#3", 0),
    ("iter-10: perfsim search rank-3 — −46% TPS (11.2× off)", 0),
])

# Left column: durable artifacts
add_text(s, Inches(0.5), Inches(4.5), Inches(6.2), Inches(0.4),
         "What landed despite no perf win:",
         size=16, bold=True, color=GREEN)
add_bullets(s, Inches(0.5), Inches(4.95), Inches(6.2), Inches(2.0), [
    ("4 issues filed (jax-gpt#2/#3, perfsim#44/#47)", 0),
    ("3 perfsim PRs merged (#45 #46 + #48 corpus backfill)", 0),
    ("2 corpus anchors (production + iter-10 calibration miss)", 0),
    ("AGENT.md self-improvements: rehydration block,", 0),
    ("Step 12.5 corpus rule, lever-source taxonomy", 1),
])

# Right column: non-cluster iters (diagnosis)
add_text(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(0.4),
         "Diagnostic iters (no cluster spend):",
         size=16, bold=True, color=BLUE)
add_bullets(s, Inches(7.0), Inches(2.55), Inches(5.8), Inches(4.0), [
    ("iter-3: BF16 microbench grid → closed perfsim#10", 0),
    ("iter-4: moe_gmm_ag bisection — 16,656 ms (48% of step)", 0),
    ("iter-6: /checkpoint/ bisection — 4,216 ms attention recompute", 0),
    ("iter-9: perfsim cross-check — predicted +4.5% upside", 0),
    ("iter-11: perfsim corpus backfill (Step 12.5 retroactive)", 0),
    ("iter-12: prior-art survey — diagnosed inapplicable", 0),
    ("iter-13: chunk-pipelining timeline — body-tail exposure", 0),
    ("iter-14: ragged_a2a AOT probe → abandon path C", 0),
])

# Box at bottom: takeaway
add_box(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
        fill=RGBColor(0xFF, 0xF4, 0xE5), border=ACCENT)
add_text(s, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
         "Cumulative HALT at iter-12 — single-iter levers on iter-2b exhausted. Multi-iter scope authorization unlocked iter-15+16.",
         size=13, color=GRAY, bold=True)

add_footer(s, "honest middle · iters 3-14 · 4 cluster shots, 0 perf gains, lots of durable knowledge")

# ── Slide 12: iter-16 attention-only-checkpoint ────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Iter-16 — attention-only-checkpoint LANDS")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "Multi-iter scope #2. One-line patch. First cluster-validated gain since iter-2 (4 days).",
         size=18, color=GRAY)

# left: results table
tbl = s.shapes.add_table(rows=5, cols=4,
                         left=Inches(0.5), top=Inches(2.1),
                         width=Inches(6.8), height=Inches(2.5)).table
hdr = ["Metric", "iter-2b", "iter-16", "Δ"]
rows = [
    ["step time", "34,659 ms", "34,200 ms", "−460 ms (−1.3%)"],
    ["TPS/chip", "1882", "1916", "+34 (+1.8%)"],
    ["MFU", "30.5%", "31.1%", "+0.6 pp"],
    ["loss (lm/aux)", "12.037/403.45", "12.037/403.43", "matches ✓"],
]
for j, txt in enumerate(hdr):
    cell = tbl.cell(0, j)
    cell.text = txt
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(13)
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
for i, row in enumerate(rows, 1):
    for j, txt in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = txt
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12)

# right: the patch + insight
add_text(s, Inches(7.6), Inches(2.1), Inches(5.5), Inches(0.4),
         "The one-line change",
         size=15, bold=True, color=BLUE)
add_text(s, Inches(7.6), Inches(2.5), Inches(5.5), Inches(0.4),
         "model.py:3053",
         size=11, color=GRAY)
add_box(s, Inches(7.6), Inches(2.85), Inches(5.5), Inches(1.4),
        fill=RGBColor(0xF5, 0xF5, 0xF5), border=GRAY)
add_text(s, Inches(7.8), Inches(2.95), Inches(5.3), Inches(0.4),
         "names_which_can_be_saved=",
         size=12, color=RGBColor(0x33, 0x33, 0x33), font="Courier New")
add_text(s, Inches(7.8), Inches(3.3), Inches(5.3), Inches(0.4),
         '   ("attn_proj_out",)',
         size=12, color=GREEN, font="Courier New", bold=True)
add_text(s, Inches(7.8), Inches(3.65), Inches(5.3), Inches(0.4),
         "# was: ()",
         size=11, color=GRAY, font="Courier New")

add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.4),
         "Why it works:",
         size=15, bold=True, color=BLUE)
add_bullets(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), [
    ("attn_proj_out was already checkpoint_named at model.py:560,636 with comment \"skip Splash bwd recompute\"", 0),
    ("The author anticipated this exact lever 4+ days before iter-16 found it", 1),
    ("Saving it in HBM (~26 GB / 58 layers) lets bwd skip the 4,216 ms attention recompute (iter-6 finding)", 0),
    ("Smaller than perfsim's +4.5% prediction (FFN intermediates still recomputed — perfsim used a half-of-none heuristic)", 0),
])
add_footer(s, "iter-16 · cluster-validated 2026-05-11 · IMPROVED candidate pending iter-17 ratchet")

# ── Slide 13: SAVE vs OFFLOAD ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "The load-bearing distinction: SAVE vs OFFLOAD")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "iter-7 NaN'd. iter-16 won. Same tensor name. One word different.",
         size=18, color=GRAY)

# Two side-by-side boxes
def policy_card(s, x, y, w, h, title, code1, code2, outcome, outcome_color):
    box = add_box(s, x, y, w, h, fill=LIGHT, border=BLUE)
    add_text(s, x + Inches(0.25), y + Inches(0.2), w - Inches(0.5), Inches(0.5),
             title, size=18, bold=True, color=BLUE)
    add_box(s, x + Inches(0.25), y + Inches(0.85), w - Inches(0.5), Inches(1.3),
            fill=RGBColor(0xF5, 0xF5, 0xF5), border=GRAY)
    add_text(s, x + Inches(0.45), y + Inches(0.95), w - Inches(0.9), Inches(0.4),
             code1, size=11, color=RGBColor(0x33, 0x33, 0x33), font="Courier New")
    add_text(s, x + Inches(0.45), y + Inches(1.3), w - Inches(0.9), Inches(0.4),
             code2, size=11, color=ACCENT, font="Courier New", bold=True)
    add_text(s, x + Inches(0.45), y + Inches(1.7), w - Inches(0.9), Inches(0.4),
             "# different code path!", size=10, color=GRAY, font="Courier New")
    add_text(s, x + Inches(0.25), y + Inches(2.4), w - Inches(0.5), Inches(0.5),
             outcome, size=15, bold=True, color=outcome_color)

policy_card(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(3.2),
            "iter-7 — OFFLOAD list",
            'names_which_can_be_offloaded=',
            '  ("attn_proj_out",)',
            "→ NaN at step 1 (jax-gpt#2)",
            ACCENT)

policy_card(s, Inches(6.75), Inches(2.0), Inches(6.1), Inches(3.2),
            "iter-16 — SAVE list",
            'names_which_can_be_saved=',
            '  ("attn_proj_out",)',
            "→ +1.8% TPS, loss matches",
            GREEN)

add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
         "Why the distinction:",
         size=15, bold=True, color=BLUE)
add_bullets(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.4), [
    ("SAVE → HBM (device) — direct save/restore. Works.", 0),
    ("OFFLOAD → pinned_host — async DMA + restore. Broken: jax-gpt#2/#3.", 0),
    ("iter-7's jax-gpt#2 read \"the lever is dead\"; only the OFFLOAD code path was dead", 0),
    ("iter-16 encoded in AGENT.md §5c + §13: NaN issues MUST enumerate untried alternative code paths", 0),
])
add_footer(s, "systematized · AGENT.md §5c \"Working with checkpoint policies\"")

# ── Slide 14: where we're headed (refreshed) ────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Where we're headed")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "Ratchet iter-16 → ratchet baseline; then the next multi-iter scope.",
         size=18, color=GRAY)

# left: near-term
add_text(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(0.4),
         "Near-term (this week)",
         size=18, bold=True, color=BLUE)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(5.8), Inches(4.0), [
    ("Iter-17 = repeat measurement on iter-16 image", 0),
    ("Confirms gain within ±0.3% noise band → promote to BASELINE", 1),
    ("Ratchet corpus production entry from 34,659 → 34,200 ms", 1),
    ("Land jax-gpt#2/#3 fix (maintainer) → unblocks OFFLOAD lever class", 0),
    ("Add cde.yaml to .dockerignore (build-hash leak fix)", 0),
    ("Update perfsim's training-regime remat model to know about", 0),
    ("attn_proj_out save (currently predicts attn_only at +4.5%)", 1),
])

# right: longer-term
add_text(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(0.4),
         "Multi-iter scope queue (need authorization)",
         size=18, bold=True, color=BLUE)
add_bullets(s, Inches(7.0), Inches(2.5), Inches(5.8), Inches(4.0), [
    ("Custom Pallas tgmm with vmem_limit_bytes (re-open iter-5)", 0),
    ("Chunk-pipelining/overlap fix (iter-13 finding: ~2.4% in body-tail)", 0),
    ("Fused gate+up GMM (needs checkpoint reformat + MMLU re-run)", 0),
    ("After upstream NaN fixes: re-attempt OFFLOAD-class levers", 0),
    ("attn_proj_out OFFLOAD = host-resident, frees ~26 GB HBM", 1),
    ("Compose with iter-16 SAVE for stacking-class gains", 1),
])

add_footer(s, "iter-17 ratchet next · multi-iter scope queue ready when authorized")

# ── Slide 15: lessons + closing ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "What we've learned")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
         "Four observations from the 16-iter arc that shape how we build agentic systems.",
         size=18, color=GRAY)

# 4 columns of lessons
def lesson_card(s, x, y, w, h, num, title, body):
    box = add_box(s, x, y, w, h, fill=LIGHT, border=BLUE)
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
             num, size=36, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.75), y + Inches(0.3), w - Inches(0.95), Inches(0.4),
             title, size=14, bold=True, color=BLUE)
    add_text(s, x + Inches(0.2), y + Inches(1.0), w - Inches(0.4), h - Inches(1.1),
             body, size=11, color=GRAY)

lesson_card(s, Inches(0.4), Inches(2.0), Inches(3.1), Inches(4.5), "1",
            "Trajectory > snapshot",
            "+1.8% TPS in iter-16 is small. The trajectory matters more: "
            "calibration unblocks search, search unblocks iteration, "
            "iteration compounds. The harness self-improvements across "
            "16 iters (corpus protocol, rehydration, lever-source taxonomy) "
            "are bigger than the perf gain.")
lesson_card(s, Inches(3.65), Inches(2.0), Inches(3.1), Inches(4.5), "2",
            "Doing vs reviewing",
            "4-agent (one doer per repo) had multi-hour handoff latency. "
            "Collapsing doing into one agent + async PR review preserved "
            "the second-pair-of-eyes safety net without the synchronous "
            "handoff cost. 12 of 13 perfsim deep-review issues closed in "
            "3 hours.")
lesson_card(s, Inches(6.9), Inches(2.0), Inches(3.1), Inches(4.5), "3",
            "GitHub as message bus",
            "No A2A protocol, no orchestrator. Issues + PRs give us free "
            "durable audit trails, async semantics, LLM-agnostic ops. "
            "Worktrees + per-branch isolation handle cross-repo work. "
            "Every cluster shot anchors the perfsim corpus.")
lesson_card(s, Inches(10.15), Inches(2.0), Inches(3.1), Inches(4.5), "4",
            "Enumerate untried code paths",
            "iter-7 NaN'd on OFFLOAD attn_proj_out and filed jax-gpt#2 as "
            "\"lever dead\". The SAVE-list variant (same name, different "
            "code path) sat undiscovered for 4 days. AGENT.md §13 now "
            "requires NaN issue bodies to enumerate untried alternative "
            "code paths.")

add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         "Questions, feedback, or want to drive an iteration? Branch is autoperf/dsv3_train_full.",
         size=14, color=BLUE, align=PP_ALIGN.CENTER, bold=True)
add_footer(s, "closing · 2026-05-11 · 16 iters · 2 cluster-validated wins · 5 perfsim PRs")

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/sivaibhav_google_com/jax-gpt/docs/autoperf-overview.pptx"
prs.save(out)
print(f"saved: {out}")
print(f"slides: {len(prs.slides)}")
